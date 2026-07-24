import os
import shutil
import glob
import re
import json
from typing import List, Dict, Any, Optional
from pathlib import Path
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings

# ============ FILE-SPECIFIC EXTRACTORS ============

def extract_text_from_ipynb(file_path: str) -> str:
    """Extract text from Jupyter Notebook files"""
    try:
        import nbformat
        from nbconvert import PythonExporter
        
        with open(file_path, 'r', encoding='utf-8') as f:
            notebook = nbformat.read(f, as_version=4)
        
        extracted_parts = []
        
        # Add notebook metadata
        if hasattr(notebook, 'metadata'):
            metadata = notebook.metadata
            if 'kernelspec' in metadata:
                extracted_parts.append(f"KERNEL: {metadata['kernelspec'].get('display_name', 'Unknown')}")
        
        cell_count = 0
        for cell in notebook.cells:
            cell_count += 1
            cell_type = cell.cell_type
            
            if cell_type == 'markdown':
                extracted_parts.append(f"CELL {cell_count} (MARKDOWN):\n{cell.source}")
                
            elif cell_type == 'code':
                code_content = cell.source
                extracted_parts.append(f"CELL {cell_count} (CODE):\n{code_content}")
                
                # Extract output text if present
                if 'outputs' in cell and cell.outputs:
                    output_texts = []
                    for output in cell.outputs:
                        if 'text' in output:
                            output_texts.append(''.join(output.text))
                        elif 'data' in output and 'text/plain' in output.data:
                            output_texts.append(output.data['text/plain'])
                    
                    if output_texts:
                        extracted_parts.append(f"OUTPUT:\n{''.join(output_texts)}")
            
            elif cell_type == 'raw':
                extracted_parts.append(f"CELL {cell_count} (RAW):\n{cell.source}")
        
        return "\n\n".join(extracted_parts)
    
    except Exception as e:
        return f"ERROR processing notebook: {str(e)}"

def extract_text_from_py(file_path: str) -> str:
    """Extract text from Python files"""
    try:
        import ast
        
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        extracted_parts = []
        
        # Extract docstrings
        try:
            tree = ast.parse(content)
            docstrings = []
            
            for node in ast.walk(tree):
                if isinstance(node, (ast.FunctionDef, ast.ClassDef, ast.Module)):
                    docstring = ast.get_docstring(node)
                    if docstring:
                        node_type = node.__class__.__name__
                        node_name = getattr(node, 'name', 'module')
                        docstrings.append(f"{node_type} '{node_name}':\n{docstring}")
            
            if docstrings:
                extracted_parts.append("DOCSTRINGS:\n" + "\n\n".join(docstrings))
        except SyntaxError:
            pass
        
        # Extract imports
        imports = []
        try:
            tree = ast.parse(content)
            for node in ast.walk(tree):
                if isinstance(node, ast.Import):
                    for alias in node.names:
                        imports.append(f"import {alias.name}")
                elif isinstance(node, ast.ImportFrom):
                    module = node.module or ''
                    for alias in node.names:
                        imports.append(f"from {module} import {alias.name}")
            
            if imports:
                extracted_parts.append("IMPORTS:\n" + "\n".join(imports[:50]))
        except:
            pass
        
        # Extract comments
        comments = re.findall(r'#.*$', content, re.MULTILINE)
        if comments:
            cleaned_comments = [c.strip() for c in comments if c.strip()]
            if cleaned_comments:
                extracted_parts.append("COMMENTS:\n" + "\n".join(cleaned_comments[:100]))
        
        # Extract code (first 200 lines)
        code_lines = []
        in_docstring = False
        docstring_delimiter = None
        
        for line in content.split('\n'):
            stripped = line.strip()
            
            if not stripped or stripped.startswith('#'):
                continue
            
            if stripped.startswith('"""') or stripped.startswith("'''"):
                if not in_docstring:
                    in_docstring = True
                    docstring_delimiter = stripped[:3]
                    if stripped.count(docstring_delimiter) >= 2:
                        in_docstring = False
                    continue
                else:
                    if stripped.endswith(docstring_delimiter):
                        in_docstring = False
                    continue
            
            if not in_docstring:
                code_lines.append(line)
        
        if code_lines:
            code_text = "\n".join(code_lines[:200])
            if len(code_lines) > 200:
                code_text += f"\n... (truncated, {len(code_lines) - 200} more lines)"
            extracted_parts.append("CODE:\n" + code_text)
        
        if not extracted_parts:
            lines = content.split('\n')[:500]
            return "FILE CONTENT:\n" + "\n".join(lines)
        
        return "\n\n".join(extracted_parts)
    
    except Exception as e:
        return f"ERROR processing Python file: {str(e)}"

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF files"""
    try:
        import PyPDF2
        
        text_parts = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Get metadata
            if pdf_reader.metadata:
                meta = pdf_reader.metadata
                if '/Title' in meta:
                    text_parts.append(f"TITLE: {meta['/Title']}")
                if '/Author' in meta:
                    text_parts.append(f"AUTHOR: {meta['/Author']}")
            
            # Extract text from each page
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"PAGE {page_num}:\n{text}")
        
        return "\n\n".join(text_parts) if text_parts else "No text extracted from PDF"
    
    except Exception as e:
        return f"ERROR processing PDF: {str(e)}"

def extract_text_from_docx(file_path: str) -> str:
    """Extract text from Word documents"""
    try:
        import docx
        
        doc = docx.Document(file_path)
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    table_text.append(" | ".join(row_text))
            if table_text:
                text_parts.append("TABLE:\n" + "\n".join(table_text))
        
        return "\n".join(text_parts) if text_parts else "No text extracted from DOCX"
    
    except Exception as e:
        return f"ERROR processing DOCX: {str(e)}"

def extract_text_from_excel(file_path: str) -> str:
    """Extract text from Excel files"""
    try:
        import openpyxl
        
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = [f"SHEET: {sheet_name}"]
            
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else '' for cell in row]
                if any(row_values):
                    sheet_text.append(" | ".join(row_values))
            
            if len(sheet_text) > 1:
                text_parts.append("\n".join(sheet_text))
        
        return "\n\n".join(text_parts) if text_parts else "No text extracted from Excel"
    
    except Exception as e:
        return f"ERROR processing Excel: {str(e)}"

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PowerPoint presentations"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"SLIDE {slide_num}:"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                elif hasattr(shape, "table"):
                    table_text = []
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(" | ".join(row_text))
                    if table_text:
                        slide_text.append("TABLE:\n" + "\n".join(table_text))
            
            if hasattr(slide, "notes_slide") and slide.notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    slide_text.append(f"NOTES:\n{notes_text}")
            
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts) if text_parts else "No text extracted from PPTX"
    
    except Exception as e:
        return f"ERROR processing PPTX: {str(e)}"

def extract_text_from_image(file_path: str) -> str:
    """Extract text from images using OCR"""
    try:
        from PIL import Image
        import pytesseract
        
        image = Image.open(file_path)
        # Convert to grayscale for better OCR
        image = image.convert('L')
        text = pytesseract.image_to_string(image)
        
        if text.strip():
            return f"IMAGE TEXT:\n{text}"
        else:
            return "No text detected in image"
    
    except Exception as e:
        return f"ERROR processing image: {str(e)}"

def extract_text_from_file(file_path: str) -> str:
    """Route file to appropriate extractor based on extension"""
    ext = os.path.splitext(file_path)[1].lower()
    
    extractors = {
        '.ipynb': extract_text_from_ipynb,
        '.py': extract_text_from_py,
        '.pdf': extract_text_from_pdf,
        '.docx': extract_text_from_docx,
        '.xlsx': extract_text_from_excel,
        '.pptx': extract_text_from_pptx,
        '.jpg': extract_text_from_image,
        '.jpeg': extract_text_from_image,
        '.png': extract_text_from_image,
        '.txt': lambda f: f"FILE CONTENT:\n{open(f, 'r', encoding='utf-8').read()}"
    }
    
    extractor = extractors.get(ext)
    if extractor:
        try:
            return extractor(file_path)
        except Exception as e:
            return f"ERROR: {str(e)}"
    else:
        return f"Unsupported file type: {ext}"

# ============ CUSTOM DOCUMENT LOADER ============

class CustomDocumentLoader:
    """Custom loader for all supported file formats"""
    
    def __init__(self, directory: str):
        self.directory = directory
        self.supported_extensions = {
            '.ipynb', '.py', '.pdf', '.docx', '.xlsx', '.pptx', 
            '.jpg', '.jpeg', '.png', '.txt'
        }
    
    def load_documents(self) -> List[Document]:
        """Load all supported documents from directory recursively"""
        documents = []
        
        if not os.path.exists(self.directory):
            print(f"Warning: {self.directory} directory not found!")
            return documents
        
        # Recursively find all supported files
        for ext in self.supported_extensions:
            pattern = os.path.join(self.directory, f"**/*{ext}")
            files = glob.glob(pattern, recursive=True)
            
            for file_path in files:
                try:
                    # Extract text content
                    text = extract_text_from_file(file_path)
                    
                    # Get relative path for metadata
                    rel_path = os.path.relpath(file_path, self.directory)
                    
                    # Create Document object
                    doc = Document(
                        page_content=text,
                        metadata={
                            'source': file_path,
                            'filename': os.path.basename(file_path),
                            'relative_path': rel_path,
                            'extension': os.path.splitext(file_path)[1].lower(),
                            'type': self._get_file_type(file_path)
                        }
                    )
                    documents.append(doc)
                    print(f"✅ Loaded: {rel_path}")
                    
                except Exception as e:
                    print(f"❌ Error loading {file_path}: {e}")
        
        return documents
    
    def _get_file_type(self, file_path: str) -> str:
        """Get file type category"""
        ext = os.path.splitext(file_path)[1].lower()
        type_map = {
            '.ipynb': 'jupyter_notebook',
            '.py': 'python',
            '.pdf': 'pdf',
            '.docx': 'word',
            '.xlsx': 'excel',
            '.pptx': 'powerpoint',
            '.jpg': 'image',
            '.jpeg': 'image',
            '.png': 'image',
            '.txt': 'text'
        }
        return type_map.get(ext, 'unknown')

# ============ MAIN PROCESSING FUNCTIONS ============

def load_documents():
    """Load all documents from the course materials directory using custom loader"""
    loader = CustomDocumentLoader(COURSE_MATERIALS_DIR)
    documents = loader.load_documents()
    return documents

def process_and_store_documents():
    """Process documents and store in vector database"""
    
    # Delete existing database to avoid dimension mismatch
    if os.path.exists(DB_DIR):
        print(f"Removing existing database at {DB_DIR}...")
        shutil.rmtree(DB_DIR)
        print("Database removed successfully.")
    
    print("Loading documents...")
    documents = load_documents()
    
    if not documents:
        print("No documents found! Please add files to the course materials directory.")
        print("Supported formats: IPYNB, PY, PDF, DOCX, XLSX, PPTX, JPG, JPEG, PNG, TXT")
        return
    
    print(f"Loaded {len(documents)} documents. Splitting into chunks...")
    
    # Split documents - using a more appropriate splitter for code and mixed content
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50,
        separators=[
            "\n\n",  # Paragraphs
            "\n",    # Lines
            ". ",    # Sentences
            " ",     # Words
            ""       # Characters
        ],
        length_function=len,
    )
    chunks = text_splitter.split_documents(documents)
    print(f"Created {len(chunks)} chunks")
    
    # Initialize embeddings
    print("Initializing embeddings...")
    embeddings = HuggingFaceEmbeddings(
        model_name="all-MiniLM-L6-v2",
        model_kwargs={'device': 'cpu'},
        encode_kwargs={'normalize_embeddings': True}
    )
    
    # Create vector database - auto-persists with persist_directory
    print("Creating vector database...")
    vector_db = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory=DB_DIR
    )
    
    print(f"✅ Vector database created successfully at {DB_DIR}")
    print(f"Total chunks stored: {len(chunks)}")

# ============ CONFIGURATION ============
DB_DIR = "./chroma_db"
COURSE_MATERIALS_DIR = "./docs"

# ============ MAIN ============
if __name__ == "__main__":
    process_and_store_documents()